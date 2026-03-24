#!/usr/bin/env python3
"""
PX Tech Report PPTX Generator
==============================
Design System baseado no Boletim Tech Jan 2026 da PX Ativos Judiciais.

Uso:
    python pptx_generator.py --type weekly --output report.pptx
    python pptx_generator.py --type general --output report.pptx --data metrics.json

Dependências:
    pip install python-pptx

Author: Cortex AI Framework
"""

import argparse
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import nsmap
except ImportError as e:
    print(f"Erro: python-pptx não instalado ou erro de importação: {e}")
    print("Execute: pip install python-pptx")
    exit(1)

# Alias for compatibility
RgbColor = RGBColor


# =============================================================================
# DESIGN SYSTEM - Cores PX
# =============================================================================

class PXColors:
    """Paleta de cores oficial PX Tech Report."""

    # Primary
    NAVY_BLUE = RgbColor(0x1E, 0x2A, 0x3A)       # #1E2A3A - Background principal
    AMBER = RgbColor(0xF5, 0xA6, 0x23)            # #F5A623 - Destaque
    WHITE = RgbColor(0xFF, 0xFF, 0xFF)            # #FFFFFF - Texto principal
    LIGHT_GRAY = RgbColor(0xA0, 0xA0, 0xA0)       # #A0A0A0 - Texto secundário

    # Cards
    CARD_NAVY = RgbColor(0x2D, 0x3E, 0x50)        # #2D3E50
    CARD_AMBER = RgbColor(0xF5, 0xA6, 0x23)       # #F5A623
    CARD_ORANGE = RgbColor(0xE6, 0x7E, 0x22)      # #E67E22
    CARD_GREEN = RgbColor(0x27, 0xAE, 0x60)       # #27AE60
    CARD_LIGHT = RgbColor(0xF5, 0xF5, 0xF5)       # #F5F5F5

    # Status
    SUCCESS = RgbColor(0x27, 0xAE, 0x60)          # #27AE60
    WARNING = RgbColor(0xF5, 0xA6, 0x23)          # #F5A623
    IN_PROGRESS = RgbColor(0xE6, 0x7E, 0x22)      # #E67E22
    INFO = RgbColor(0x34, 0x98, 0xDB)             # #3498DB


# =============================================================================
# CONFIGURAÇÃO DE SLIDES
# =============================================================================

class SlideConfig:
    """Configurações de layout."""

    # Dimensões (16:9)
    WIDTH = Inches(13.333)
    HEIGHT = Inches(7.5)

    # Margens
    MARGIN_LEFT = Inches(0.6)
    MARGIN_RIGHT = Inches(0.6)
    MARGIN_TOP = Inches(0.8)
    MARGIN_BOTTOM = Inches(0.5)

    # Header strip
    HEADER_HEIGHT = Inches(0.6)
    HEADER_ACCENT_SIZE = Inches(0.4)

    # Footer
    FOOTER_HEIGHT = Inches(0.4)

    # Fonts
    FONT_TITLE = "Inter"
    FONT_BODY = "Inter"

    # Font sizes
    SIZE_TITLE = Pt(48)
    SIZE_SUBTITLE = Pt(24)
    SIZE_METRIC = Pt(72)
    SIZE_BODY = Pt(18)
    SIZE_CAPTION = Pt(14)
    SIZE_FOOTER = Pt(10)


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_shape_fill(shape, color: RgbColor):
    """Define cor de preenchimento de uma shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_strip(slide, prs):
    """Adiciona header strip navy com ícone amber no canto."""
    # Header strip navy
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, SlideConfig.HEADER_HEIGHT
    )
    set_shape_fill(header, PXColors.NAVY_BLUE)
    header.line.fill.background()

    # Ícone amber no canto direito
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        prs.slide_width - Inches(0.8), Inches(0.1),
        SlideConfig.HEADER_ACCENT_SIZE, SlideConfig.HEADER_ACCENT_SIZE
    )
    set_shape_fill(accent, PXColors.AMBER)
    accent.line.fill.background()


def add_footer(slide, prs, text: str, page_num: int):
    """Adiciona footer com texto e número da página."""
    # Footer text
    footer_left = slide.shapes.add_textbox(
        SlideConfig.MARGIN_LEFT,
        prs.slide_height - SlideConfig.FOOTER_HEIGHT,
        Inches(6), SlideConfig.FOOTER_HEIGHT
    )
    tf = footer_left.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = SlideConfig.SIZE_FOOTER
    p.font.color.rgb = PXColors.LIGHT_GRAY
    p.font.name = SlideConfig.FONT_BODY

    # Page number
    footer_right = slide.shapes.add_textbox(
        prs.slide_width - Inches(1),
        prs.slide_height - SlideConfig.FOOTER_HEIGHT,
        Inches(0.5), SlideConfig.FOOTER_HEIGHT
    )
    tf = footer_right.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num:02d}"
    p.font.size = SlideConfig.SIZE_FOOTER
    p.font.color.rgb = PXColors.LIGHT_GRAY
    p.font.name = SlideConfig.FONT_BODY
    p.alignment = PP_ALIGN.RIGHT


def add_title_with_highlight(slide, title: str, highlight_word: str,
                             top: float = 1.2, font_size: Pt = None):
    """Adiciona título com palavra destacada em amber."""
    font_size = font_size or SlideConfig.SIZE_TITLE

    title_box = slide.shapes.add_textbox(
        SlideConfig.MARGIN_LEFT, Inches(top),
        Inches(10), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]

    # Split title em partes
    parts = title.split(highlight_word)
    if len(parts) == 2:
        # Parte antes do highlight
        if parts[0]:
            run = p.add_run()
            run.text = parts[0]
            run.font.size = font_size
            run.font.color.rgb = PXColors.WHITE
            run.font.name = SlideConfig.FONT_TITLE
            run.font.bold = True

        # Highlight
        run = p.add_run()
        run.text = highlight_word
        run.font.size = font_size
        run.font.color.rgb = PXColors.AMBER
        run.font.name = SlideConfig.FONT_TITLE
        run.font.bold = True

        # Parte depois do highlight
        if parts[1]:
            run = p.add_run()
            run.text = parts[1]
            run.font.size = font_size
            run.font.color.rgb = PXColors.WHITE
            run.font.name = SlideConfig.FONT_TITLE
            run.font.bold = True
    else:
        # Sem highlight, título normal
        p.text = title
        p.font.size = font_size
        p.font.color.rgb = PXColors.WHITE
        p.font.name = SlideConfig.FONT_TITLE
        p.font.bold = True


def add_subtitle(slide, text: str, top: float = 2.0):
    """Adiciona subtítulo em cinza."""
    subtitle_box = slide.shapes.add_textbox(
        SlideConfig.MARGIN_LEFT, Inches(top),
        Inches(10), Inches(0.5)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = SlideConfig.SIZE_SUBTITLE
    p.font.color.rgb = PXColors.LIGHT_GRAY
    p.font.name = SlideConfig.FONT_BODY


def add_metric_card(slide, left: float, top: float, width: float, height: float,
                    value: str, label: str, color: RgbColor = None):
    """Adiciona card com métrica grande."""
    color = color or PXColors.CARD_NAVY

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    set_shape_fill(card, color)
    card.line.fill.background()

    # Valor grande
    value_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.3),
        Inches(width - 0.4), Inches(height * 0.6)
    )
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = SlideConfig.SIZE_METRIC
    p.font.color.rgb = PXColors.WHITE
    p.font.name = SlideConfig.FONT_TITLE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + height - 0.6),
        Inches(width - 0.4), Inches(0.4)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = SlideConfig.SIZE_CAPTION
    p.font.color.rgb = PXColors.WHITE if color != PXColors.CARD_LIGHT else PXColors.NAVY_BLUE
    p.font.name = SlideConfig.FONT_BODY
    p.alignment = PP_ALIGN.CENTER


def add_highlight_item(slide, left: float, top: float, title: str, description: str):
    """Adiciona item com barra amber à esquerda."""
    # Barra amber
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(0.05), Inches(0.6)
    )
    set_shape_fill(bar, PXColors.AMBER)
    bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top),
        Inches(4), Inches(0.3)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = PXColors.NAVY_BLUE
    p.font.name = SlideConfig.FONT_BODY
    p.font.bold = True

    # Descrição
    desc_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.25),
        Inches(4), Inches(0.3)
    )
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(12)
    p.font.color.rgb = PXColors.LIGHT_GRAY
    p.font.name = SlideConfig.FONT_BODY


# =============================================================================
# SLIDE GENERATORS
# =============================================================================

def create_cover_slide(prs, data: Dict) -> None:
    """Cria slide de capa."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background navy
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    set_shape_fill(background, PXColors.NAVY_BLUE)
    background.line.fill.background()

    # Ícone amber no canto
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        prs.slide_width - Inches(0.8), Inches(0.4),
        Inches(0.5), Inches(0.5)
    )
    set_shape_fill(accent, PXColors.AMBER)
    accent.line.fill.background()

    # Pill amarelo (tag)
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SlideConfig.MARGIN_LEFT, Inches(1.8),
        Inches(2), Inches(0.4)
    )
    set_shape_fill(pill, PXColors.AMBER)
    pill.line.fill.background()

    # Título
    add_title_with_highlight(
        slide,
        data.get('title', 'TECH REPORT 2025'),
        data.get('highlight', 'REPORT'),
        top=2.5,
        font_size=Pt(56)
    )

    # Subtítulo
    add_subtitle(slide, data.get('subtitle', 'Consolidado de entregas tecnológicas'), top=3.8)

    # Métricas na parte inferior
    metrics = data.get('cover_metrics', [
        {'label': 'Projetos entregues', 'value': '15+'},
        {'label': 'Áreas impactadas', 'value': '8'},
        {'label': 'AI First', 'value': '100%'}
    ])

    x_start = 0.6
    for i, metric in enumerate(metrics[:3]):
        metric_box = slide.shapes.add_textbox(
            Inches(x_start + i * 2.5), Inches(5.5),
            Inches(2), Inches(0.8)
        )
        tf = metric_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['label']
        p.font.size = Pt(12)
        p.font.color.rgb = PXColors.LIGHT_GRAY
        p.font.name = SlideConfig.FONT_BODY

    # Footer
    footer_box = slide.shapes.add_textbox(
        SlideConfig.MARGIN_LEFT, Inches(6.8),
        Inches(4), Inches(0.3)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('footer', 'pxativosjudiciais.com.br')
    p.font.size = Pt(10)
    p.font.color.rgb = PXColors.LIGHT_GRAY
    p.font.name = SlideConfig.FONT_BODY


def create_summary_slide(prs, data: Dict, page_num: int) -> None:
    """Cria slide de resumo com métricas."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background branco
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    set_shape_fill(background, PXColors.WHITE)
    background.line.fill.background()

    add_header_strip(slide, prs)

    # Título
    add_title_with_highlight(
        slide,
        data.get('title', 'Um ano de transformação digital'),
        data.get('highlight', 'transformação digital'),
        top=1.0,
        font_size=Pt(40)
    )
    add_subtitle(slide, data.get('subtitle', 'Principais conquistas e marcos alcançados'), top=1.7)

    # Métricas cards
    metrics = data.get('metrics', [
        {'value': '500+', 'label': 'Tickets atendidos', 'color': 'navy'},
        {'value': '200+', 'label': 'Chamados resolvidos', 'color': 'amber'},
        {'value': '100%', 'label': 'Nova infraestrutura', 'color': 'light'},
        {'value': '6', 'label': 'Projetos com IA', 'color': 'light'}
    ])

    colors_map = {
        'navy': PXColors.CARD_NAVY,
        'amber': PXColors.CARD_AMBER,
        'orange': PXColors.CARD_ORANGE,
        'light': PXColors.CARD_LIGHT,
        'green': PXColors.CARD_GREEN
    }

    x_start = 0.6
    for i, metric in enumerate(metrics[:4]):
        color = colors_map.get(metric.get('color', 'light'), PXColors.CARD_LIGHT)
        add_metric_card(
            slide,
            left=x_start + i * 3.1,
            top=2.4,
            width=2.8,
            height=1.8,
            value=metric['value'],
            label=metric['label'],
            color=color
        )

    # Highlights inferiores
    highlights = data.get('highlights', [
        {'title': 'Metodologia Cortex', 'description': 'Desenvolvimento AI-First acelerado'},
        {'title': 'Economia de R$ 150k', 'description': 'Comparado ao desenvolvimento tradicional'},
        {'title': 'Produtividade 4x', 'description': 'Multiplicador de eficiência'}
    ])

    x_start = 0.6
    for i, item in enumerate(highlights[:3]):
        add_highlight_item(
            slide,
            left=x_start + i * 4.2,
            top=5.0,
            title=item['title'],
            description=item['description']
        )

    add_footer(slide, prs, data.get('footer_text', 'Tech Report | PX Ativos Judiciais'), page_num)


def create_content_slide(prs, data: Dict, page_num: int) -> None:
    """Cria slide de conteúdo genérico."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background branco
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    set_shape_fill(background, PXColors.WHITE)
    background.line.fill.background()

    add_header_strip(slide, prs)

    # Ícone do módulo
    icon = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SlideConfig.MARGIN_LEFT, Inches(0.9),
        Inches(0.8), Inches(0.8)
    )
    set_shape_fill(icon, PXColors.AMBER)
    icon.line.fill.background()

    # Título do módulo
    title_box = slide.shapes.add_textbox(
        Inches(1.6), Inches(0.9),
        Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', 'Título do Módulo')
    p.font.size = Pt(32)
    p.font.color.rgb = PXColors.NAVY_BLUE
    p.font.name = SlideConfig.FONT_TITLE
    p.font.bold = True

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(
        Inches(1.6), Inches(1.45),
        Inches(8), Inches(0.4)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('subtitle', 'Descrição do módulo')
    p.font.size = Pt(16)
    p.font.color.rgb = PXColors.LIGHT_GRAY
    p.font.name = SlideConfig.FONT_BODY

    # Items de conteúdo (lado esquerdo)
    items = data.get('items', [])
    y_start = 2.3
    for i, item in enumerate(items[:4]):
        # Card de item
        item_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SlideConfig.MARGIN_LEFT, Inches(y_start + i * 0.9),
            Inches(6.5), Inches(0.75)
        )
        set_shape_fill(item_card, PXColors.CARD_LIGHT)
        item_card.line.fill.background()

        # Ícone colorido
        status_colors = {
            'success': PXColors.CARD_GREEN,
            'warning': PXColors.CARD_AMBER,
            'in_progress': PXColors.CARD_ORANGE,
            'pending': PXColors.LIGHT_GRAY
        }
        icon_color = status_colors.get(item.get('status', 'pending'), PXColors.LIGHT_GRAY)

        icon_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_start + i * 0.9 + 0.15),
            Inches(0.45), Inches(0.45)
        )
        set_shape_fill(icon_shape, icon_color)
        icon_shape.line.fill.background()

        # Texto do item
        item_box = slide.shapes.add_textbox(
            Inches(1.4), Inches(y_start + i * 0.9 + 0.1),
            Inches(5.5), Inches(0.6)
        )
        tf = item_box.text_frame

        p = tf.paragraphs[0]
        p.text = item.get('title', '')
        p.font.size = Pt(14)
        p.font.color.rgb = PXColors.NAVY_BLUE
        p.font.name = SlideConfig.FONT_BODY
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = item.get('description', '')
        p.font.size = Pt(11)
        p.font.color.rgb = PXColors.LIGHT_GRAY
        p.font.name = SlideConfig.FONT_BODY

    # Métrica destacada (lado direito)
    if 'main_metric' in data:
        metric_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.5), Inches(2.3),
            Inches(5), Inches(2)
        )
        set_shape_fill(metric_card, PXColors.CARD_NAVY)
        metric_card.line.fill.background()

        # Valor
        value_box = slide.shapes.add_textbox(
            Inches(7.7), Inches(2.5),
            Inches(4.6), Inches(1.2)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['main_metric'].get('value', '0')
        p.font.size = Pt(64)
        p.font.color.rgb = PXColors.AMBER
        p.font.name = SlideConfig.FONT_TITLE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(7.7), Inches(3.7),
            Inches(4.6), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['main_metric'].get('label', '')
        p.font.size = Pt(14)
        p.font.color.rgb = PXColors.WHITE
        p.font.name = SlideConfig.FONT_BODY
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, prs, data.get('footer_text', 'Tech Report | PX Ativos Judiciais'), page_num)


def create_thanks_slide(prs, data: Dict) -> None:
    """Cria slide de agradecimento/encerramento."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background navy
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    set_shape_fill(background, PXColors.NAVY_BLUE)
    background.line.fill.background()

    # Ícone amber
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        prs.slide_width - Inches(0.8), Inches(0.4),
        Inches(0.5), Inches(0.5)
    )
    set_shape_fill(accent, PXColors.AMBER)
    accent.line.fill.background()

    # Título "Obrigado!"
    title_box = slide.shapes.add_textbox(
        SlideConfig.MARGIN_LEFT, Inches(1.5),
        Inches(8), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Obrigado!"
    p.font.size = Pt(56)
    p.font.color.rgb = PXColors.AMBER
    p.font.name = SlideConfig.FONT_TITLE
    p.font.bold = True

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(
        SlideConfig.MARGIN_LEFT, Inches(2.5),
        Inches(10), Inches(0.5)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('subtitle', 'Continuamos nossa jornada de transformação digital')
    p.font.size = Pt(20)
    p.font.color.rgb = PXColors.WHITE
    p.font.name = SlideConfig.FONT_BODY

    # Métricas finais
    metrics = data.get('final_metrics', [
        {'value': '15+', 'label': 'Projetos entregues'},
        {'value': '8', 'label': 'Entregas Q1/26'},
        {'value': '100%', 'label': 'AI First'}
    ])

    x_start = 0.6
    for i, metric in enumerate(metrics[:3]):
        # Separador vertical
        if i > 0:
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_start + i * 3 - 0.3), Inches(3.4),
                Inches(0.02), Inches(1)
            )
            set_shape_fill(sep, PXColors.LIGHT_GRAY)
            sep.line.fill.background()

        # Valor
        value_box = slide.shapes.add_textbox(
            Inches(x_start + i * 3), Inches(3.4),
            Inches(2.5), Inches(0.8)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['value']
        p.font.size = Pt(48)
        p.font.color.rgb = PXColors.AMBER
        p.font.name = SlideConfig.FONT_TITLE
        p.font.bold = True

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x_start + i * 3), Inches(4.3),
            Inches(2.5), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['label']
        p.font.size = Pt(12)
        p.font.color.rgb = PXColors.LIGHT_GRAY
        p.font.name = SlideConfig.FONT_BODY

    # Quote box
    quote_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SlideConfig.MARGIN_LEFT, Inches(5.2),
        Inches(9), Inches(1)
    )
    set_shape_fill(quote_box, RgbColor(0x3D, 0x4E, 0x5A))
    quote_box.line.fill.background()

    # Barra amber no quote
    quote_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SlideConfig.MARGIN_LEFT, Inches(5.2),
        Inches(0.08), Inches(1)
    )
    set_shape_fill(quote_bar, PXColors.AMBER)
    quote_bar.line.fill.background()

    # Texto do quote
    quote_text = slide.shapes.add_textbox(
        Inches(0.9), Inches(5.4),
        Inches(8), Inches(0.6)
    )
    tf = quote_text.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('quote', 'Continuamos nossa jornada de transformação digital com foco em automação, inteligência artificial e eficiência operacional.')
    p.font.size = Pt(14)
    p.font.color.rgb = PXColors.WHITE
    p.font.name = SlideConfig.FONT_BODY

    # Footer
    footer_box = slide.shapes.add_textbox(
        SlideConfig.MARGIN_LEFT, Inches(6.8),
        Inches(8), Inches(0.3)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('footer', 'pxativosjudiciais.com.br | Tech Report Janeiro 2026')
    p.font.size = Pt(10)
    p.font.color.rgb = PXColors.LIGHT_GRAY
    p.font.name = SlideConfig.FONT_BODY


# =============================================================================
# MAIN GENERATOR
# =============================================================================

def generate_weekly_report(output_path: str, data: Dict = None) -> None:
    """Gera relatório semanal."""
    data = data or {}

    prs = Presentation()
    prs.slide_width = SlideConfig.WIDTH
    prs.slide_height = SlideConfig.HEIGHT

    # Slide 1: Capa
    create_cover_slide(prs, {
        'title': f"TECH REPORT SEMANA {data.get('week_number', 'XX')}",
        'highlight': 'REPORT',
        'subtitle': f"Período: {data.get('period', 'DD/MM a DD/MM/YYYY')}",
        'cover_metrics': data.get('cover_metrics', [
            {'label': 'Commits', 'value': data.get('commits', '0')},
            {'label': 'Features', 'value': data.get('features', '0')},
            {'label': 'Economia', 'value': data.get('savings', 'R$ 0')}
        ]),
        'footer': 'pxativosjudiciais.com.br'
    })

    # Slide 2: Resumo
    create_summary_slide(prs, {
        'title': f"Uma semana de {data.get('theme', 'progresso')}",
        'highlight': data.get('theme', 'progresso'),
        'subtitle': 'Principais entregas e marcos da semana',
        'metrics': data.get('metrics', [
            {'value': data.get('commits', '0'), 'label': 'Commits', 'color': 'navy'},
            {'value': data.get('features', '0'), 'label': 'Features', 'color': 'amber'},
            {'value': data.get('savings', 'R$ 0'), 'label': 'Economia', 'color': 'light'},
            {'value': data.get('hours_saved', '0h'), 'label': 'Horas economizadas', 'color': 'light'}
        ]),
        'highlights': data.get('highlights', []),
        'footer_text': 'Tech Report Semanal | PX Ativos Judiciais'
    }, page_num=1)

    # Slide 3: Entregas
    create_content_slide(prs, {
        'title': 'Entregas da Semana',
        'subtitle': 'Features e melhorias implementadas',
        'items': data.get('deliveries', []),
        'main_metric': {
            'value': data.get('total_savings', 'R$ 0'),
            'label': 'Economia total da semana'
        },
        'footer_text': 'Tech Report Semanal | PX Ativos Judiciais'
    }, page_num=2)

    # Slide 4: Contributors
    create_content_slide(prs, {
        'title': 'Contributors',
        'subtitle': 'Time envolvido nesta semana',
        'items': data.get('contributors', []),
        'main_metric': {
            'value': str(len(data.get('contributors', []))),
            'label': 'Colaboradores ativos'
        },
        'footer_text': 'Tech Report Semanal | PX Ativos Judiciais'
    }, page_num=3)

    # Slide 5: Obrigado
    create_thanks_slide(prs, {
        'subtitle': data.get('closing_message', 'Continuamos nossa jornada de transformação digital'),
        'final_metrics': data.get('final_metrics', [
            {'value': data.get('commits', '0'), 'label': 'Commits'},
            {'value': data.get('features', '0'), 'label': 'Features'},
            {'value': data.get('savings', 'R$ 0'), 'label': 'Economia'}
        ]),
        'quote': data.get('quote', 'Metodologia Cortex: 4x mais produtividade com desenvolvimento AI-First'),
        'footer': f"pxativosjudiciais.com.br | Tech Report Semana {data.get('week_number', 'XX')}"
    })

    prs.save(output_path)
    print(f"Relatório semanal gerado: {output_path}")


def generate_general_report(output_path: str, data: Dict = None) -> None:
    """Gera relatório geral do projeto."""
    data = data or {}

    prs = Presentation()
    prs.slide_width = SlideConfig.WIDTH
    prs.slide_height = SlideConfig.HEIGHT

    project_name = data.get('project_name', 'RHILO')
    year = data.get('year', datetime.now().year)

    # Slide 1: Capa
    create_cover_slide(prs, {
        'title': f"TECH REPORT {year}",
        'highlight': 'REPORT',
        'subtitle': f"Consolidado de entregas tecnológicas - {project_name}",
        'cover_metrics': [
            {'label': 'Features entregues', 'value': data.get('total_features', '0')},
            {'label': 'Economia total', 'value': data.get('total_savings', 'R$ 0')},
            {'label': 'Conclusão', 'value': data.get('completion', '0%')}
        ],
        'footer': 'pxativosjudiciais.com.br'
    })

    # Slide 2: Visão Geral
    create_summary_slide(prs, {
        'title': f"O que é o {project_name}?",
        'highlight': project_name,
        'subtitle': data.get('description', 'Plataforma legal-tech de transformação digital'),
        'metrics': data.get('overview_metrics', [
            {'value': data.get('commits_total', '0'), 'label': 'Commits totais', 'color': 'navy'},
            {'value': data.get('contributors_total', '0'), 'label': 'Contributors', 'color': 'amber'},
            {'value': data.get('months_dev', '0'), 'label': 'Meses de desenvolvimento', 'color': 'light'},
            {'value': data.get('adrs', '0'), 'label': 'ADRs documentados', 'color': 'light'}
        ]),
        'highlights': data.get('tech_stack', [
            {'title': 'Backend', 'description': 'NestJS + Prisma + PostgreSQL'},
            {'title': 'Frontend', 'description': 'React + TypeScript + TailwindCSS'},
            {'title': 'Infra', 'description': 'AWS + Docker + GitHub Actions'}
        ]),
        'footer_text': f'Tech Report {year} | PX Ativos Judiciais'
    }, page_num=1)

    # Slide 3: Metodologia Cortex
    create_content_slide(prs, {
        'title': 'Metodologia Cortex',
        'subtitle': 'Desenvolvimento AI-First',
        'items': [
            {'title': 'Claude Code + Agentes Especializados', 'description': 'Framework de desenvolvimento com IA', 'status': 'success'},
            {'title': 'Multiplicador 4x', 'description': 'Produtividade comparada ao tradicional', 'status': 'success'},
            {'title': 'Documentação Automática', 'description': 'ADRs e specs gerados com IA', 'status': 'success'},
            {'title': 'Testes Assistidos', 'description': 'Cobertura de testes com suporte de IA', 'status': 'success'}
        ],
        'main_metric': {
            'value': '4x',
            'label': 'Multiplicador de produtividade'
        },
        'footer_text': f'Tech Report {year} | PX Ativos Judiciais'
    }, page_num=2)

    # Slide 4: ROI Total
    create_content_slide(prs, {
        'title': 'Economia Total',
        'subtitle': 'ROI do desenvolvimento com Cortex',
        'items': data.get('roi_breakdown', [
            {'title': 'Horas economizadas', 'description': data.get('hours_saved', '0 horas'), 'status': 'success'},
            {'title': 'Custo economizado', 'description': data.get('cost_saved', 'R$ 0'), 'status': 'success'},
            {'title': 'Equivalente tradicional', 'description': data.get('dev_equivalent', '0 desenvolvedores'), 'status': 'in_progress'},
            {'title': 'Time real', 'description': data.get('actual_team', '0 pessoas'), 'status': 'success'}
        ]),
        'main_metric': {
            'value': data.get('total_savings', 'R$ 0'),
            'label': 'Economia total acumulada'
        },
        'footer_text': f'Tech Report {year} | PX Ativos Judiciais'
    }, page_num=3)

    # Slide 5: Módulos
    create_content_slide(prs, {
        'title': 'Módulos Entregues',
        'subtitle': 'Componentes principais do sistema',
        'items': data.get('modules', []),
        'main_metric': {
            'value': str(len(data.get('modules', []))),
            'label': 'Módulos implementados'
        },
        'footer_text': f'Tech Report {year} | PX Ativos Judiciais'
    }, page_num=4)

    # Slide 6: Obrigado
    create_thanks_slide(prs, {
        'subtitle': f'{year} foi um ano de grandes conquistas tecnológicas',
        'final_metrics': [
            {'value': data.get('total_features', '0'), 'label': 'Features entregues'},
            {'value': data.get('modules_count', '0'), 'label': 'Módulos implementados'},
            {'value': data.get('total_savings', 'R$ 0'), 'label': 'Economia total'}
        ],
        'quote': data.get('closing_quote', 'Metodologia Cortex: Transformando desenvolvimento de software com IA'),
        'footer': f'pxativosjudiciais.com.br | Tech Report {year}'
    })

    prs.save(output_path)
    print(f"Relatório geral gerado: {output_path}")


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description='PX Tech Report PPTX Generator',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Exemplos:
  python pptx_generator.py --type weekly --output weekly_report.pptx
  python pptx_generator.py --type general --output project_report.pptx --data metrics.json
        """
    )

    parser.add_argument(
        '--type', '-t',
        choices=['weekly', 'general'],
        required=True,
        help='Tipo de relatório a gerar'
    )

    parser.add_argument(
        '--output', '-o',
        default='report.pptx',
        help='Arquivo de saída (default: report.pptx)'
    )

    parser.add_argument(
        '--data', '-d',
        help='Arquivo JSON com dados do relatório'
    )

    args = parser.parse_args()

    # Carregar dados se fornecidos
    data = {}
    if args.data and Path(args.data).exists():
        with open(args.data, 'r', encoding='utf-8') as f:
            data = json.load(f)

    # Gerar relatório
    if args.type == 'weekly':
        generate_weekly_report(args.output, data)
    else:
        generate_general_report(args.output, data)


if __name__ == '__main__':
    main()
